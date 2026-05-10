#!/usr/bin/env python3
"""
Create PPTX presentation for the science lesson example.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

OUTPUT_PATH = "/home/z/my-project/public/downloads/cell-lesson-presentation.pptx"

# Colors
BG_DARK = RGBColor(0x0a, 0x0a, 0x0a)
BG_CARD = RGBColor(0x1a, 0x1a, 0x1a)
AMBER = RGBColor(0xf5, 0x9e, 0x0b)
GREEN = RGBColor(0x22, 0xc5, 0x5e)
PURPLE = RGBColor(0x8b, 0x5c, 0xf6)
RED = RGBColor(0xef, 0x44, 0x44)
WHITE = RGBColor(0xff, 0xff, 0xff)
GRAY = RGBColor(0x9c, 0xa3, 0xaf)
GRAY_LIGHT = RGBColor(0xd1, 0xd5, 0xdb)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Use blank layout
blank_layout = prs.slide_layouts[6]  # Blank layout

def add_bg(slide):
    """Add dark background to slide"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK

def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False, color=WHITE, align=PP_ALIGN.RIGHT, font_name='Calibri'):
    """Add a text box to the slide"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_shape_card(slide, left, top, width, height, fill_color=BG_CARD, border_color=None):
    """Add a rounded rectangle card"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_footer(slide, slide_num, total=6):
    """Add footer to slide"""
    # Footer line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(6.7), Inches(12.333), Pt(1)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0x33, 0x33, 0x33)
    line.line.fill.background()
    
    add_textbox(slide, 0.5, 6.8, 6, 0.4, "نواف البوسطه | كورس الذكاء الاصطناعي بالتعليم", 
                font_size=10, color=GRAY)
    add_textbox(slide, 6.5, 6.8, 6, 0.4, "كل المحتوى من الكتاب المدرسي وموثق بأرقام الصفحات", 
                font_size=10, color=GRAY, align=PP_ALIGN.LEFT)
    add_textbox(slide, 11.5, 0.3, 1.5, 0.4, f"{slide_num} / {total}", 
                font_size=11, color=GRAY, align=PP_ALIGN.LEFT)


# ===== SLIDE 1: Title =====
slide1 = prs.slides.add_slide(blank_layout)
add_bg(slide1)

add_textbox(slide1, 2, 1.2, 9.333, 0.5, "مثال عملي - أدوات الذكاء الاصطناعي", 
            font_size=13, color=AMBER, bold=True)

add_textbox(slide1, 2, 2, 9.333, 1.2, "الخلية الحية", 
            font_size=52, color=AMBER, bold=True)

add_textbox(slide1, 2, 3.3, 9.333, 0.6, "تركيب الخلية ومكوناتها ووظائفها", 
            font_size=26, color=GRAY)

add_textbox(slide1, 2, 4.2, 9.333, 0.5, "علوم - الصف السابع | المنهج السوري", 
            font_size=16, color=GRAY_LIGHT)

# Stats
add_textbox(slide1, 2, 5.2, 3, 0.5, "5 دقائق", font_size=28, color=AMBER, bold=True)
add_textbox(slide1, 2, 5.7, 3, 0.4, "وقت التحضير بالذكاء الاصطناعي", font_size=12, color=GRAY)

add_textbox(slide1, 5.5, 5.2, 3, 0.5, "100%", font_size=28, color=GREEN, bold=True)
add_textbox(slide1, 5.5, 5.7, 3, 0.4, "من الكتاب المدرسي", font_size=12, color=GRAY)

add_textbox(slide1, 9, 5.2, 3, 0.5, "PDF جاهز", font_size=28, color=PURPLE, bold=True)
add_textbox(slide1, 9, 5.7, 3, 0.4, "جاهز للطباعة والتوزيع", font_size=12, color=GRAY)

add_textbox(slide1, 2, 6.2, 9.333, 0.4, 
            "هاد مثال حقيقي لدرس علوم تم إعداده بالكامل باستخدام أداتي الذكاء الاصطناعي", 
            font_size=13, color=GRAY)

add_footer(slide1, 1)


# ===== SLIDE 2: Cell Structure =====
slide2 = prs.slides.add_slide(blank_layout)
add_bg(slide2)

add_textbox(slide2, 0.5, 0.3, 6, 0.4, "مثال عملي - أدوات الذكاء الاصطناعي", 
            font_size=13, color=AMBER, bold=True)

add_textbox(slide2, 0.5, 0.9, 8, 0.6, "تركيب الخلية ومكوناتها", 
            font_size=30, color=AMBER, bold=True)

add_textbox(slide2, 0.5, 1.5, 8, 0.4, "الكتاب المدرسي - صفحة 42-45", 
            font_size=16, color=GRAY)

# Left side - component cards
# Card 1: Cell Membrane
card1 = add_shape_card(slide2, 0.5, 2.2, 5.5, 1.2)
add_textbox(slide2, 0.7, 2.3, 5, 0.4, "الغشاء الخلوي", font_size=17, color=AMBER, bold=True)
add_textbox(slide2, 0.7, 2.7, 5, 0.6, "الغشاء الخارجي اللي بيحيط بالخلية وبيتحكم بالمواد اللي بتدخل وتطلع - يشبه حارس البوابة", 
            font_size=14, color=GRAY_LIGHT)

# Card 2: Cytoplasm
card2 = add_shape_card(slide2, 0.5, 3.6, 5.5, 1.2)
add_textbox(slide2, 0.7, 3.7, 5, 0.4, "السيتوبلازم", font_size=17, color=GREEN, bold=True)
add_textbox(slide2, 0.7, 4.1, 5, 0.6, "المادة الهلامية اللي بتملا الخلية وفيها العضيات - مثل السائل اللي بتحيا فيه المكونات", 
            font_size=14, color=GRAY_LIGHT)

# Card 3: Nucleus
card3 = add_shape_card(slide2, 0.5, 5.0, 5.5, 1.2)
add_textbox(slide2, 0.7, 5.1, 5, 0.4, "النواة", font_size=17, color=PURPLE, bold=True)
add_textbox(slide2, 0.7, 5.5, 5, 0.6, "مركز القيادة! فيها الحمض النووي DNA اللي بيحدد صفات الخلية وبيتحكم بنشاطاتها", 
            font_size=14, color=GRAY_LIGHT)

# Right side - Cell diagram description
card_diag = add_shape_card(slide2, 6.5, 2.2, 6.3, 4.0)
add_textbox(slide2, 7, 2.5, 5.3, 0.4, "رسم توضيحي لتركيب الخلية", font_size=16, color=WHITE, bold=True)

# Concentric circles to represent cell structure
outer = slide2.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8), Inches(3), Inches(3.3), Inches(3.3))
outer.fill.background()
outer.line.color.rgb = AMBER
outer.line.width = Pt(3)

middle = slide2.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.4), Inches(3.4), Inches(2.5), Inches(2.5))
middle.fill.background()
middle.line.color.rgb = GREEN
middle.line.width = Pt(2)
middle.line.dash_style = 4  # Dash

inner = slide2.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.1), Inches(4.1), Inches(1.1), Inches(1.1))
inner.fill.solid()
inner.fill.fore_color.rgb = RGBColor(0x2a, 0x1a, 0x4a)
inner.line.color.rgb = PURPLE
inner.line.width = Pt(3)

# Labels for diagram
add_textbox(slide2, 8.1, 2.7, 2, 0.3, "الغشاء الخلوي", font_size=11, color=AMBER, bold=True)
add_textbox(slide2, 8.5, 3.2, 2, 0.3, "السيتوبلازم", font_size=11, color=GREEN, bold=True)
add_textbox(slide2, 9.2, 4.4, 1, 0.3, "النواة", font_size=12, color=PURPLE, bold=True, align=PP_ALIGN.CENTER)

add_footer(slide2, 2)


# ===== SLIDE 3: Comparison Table =====
slide3 = prs.slides.add_slide(blank_layout)
add_bg(slide3)

add_textbox(slide3, 0.5, 0.3, 6, 0.4, "مثال عملي - أدوات الذكاء الاصطناعي", 
            font_size=13, color=AMBER, bold=True)

add_textbox(slide3, 0.5, 0.9, 10, 0.6, "مقارنة: الخلية النباتية vs الخلية الحيوانية", 
            font_size=30, color=AMBER, bold=True)

add_textbox(slide3, 0.5, 1.5, 8, 0.4, "الكتاب المدرسي - صفحة 48-50", 
            font_size=16, color=GRAY)

# Table
rows, cols = 7, 3
table = slide3.shapes.add_table(rows, cols, Inches(0.5), Inches(2.1), Inches(12.333), Inches(4.0)).table

# Set column widths
table.columns[0].width = Inches(3)
table.columns[1].width = Inches(4.666)
table.columns[2].width = Inches(4.666)

# Header row
for j, header in enumerate(["المميزة", "الخلية النباتية", "الخلية الحيوانية"]):
    cell = table.cell(0, j)
    cell.text = header
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = AMBER
        p.alignment = PP_ALIGN.CENTER
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(0x1a, 0x15, 0x0a)

# Data rows
data = [
    ("الشكل", "مستطيلة الشكل", "دائرية أو غير منتظمة"),
    ("الجدار الخلوي", "موجود", "غير موجود"),
    ("البلاستيدات الخضراء", "موجودة (للتركيب الضوئي)", "غير موجودة"),
    ("الفجوة العصارية", "كبيرة ومركزية", "صغيرة ومتعددة"),
    ("النواة", "موجودة", "موجودة"),
    ("الغشاء الخلوي", "موجود (تحت الجدار)", "موجود (الحدود الخارجية)"),
]

for i, (feat, plant, animal) in enumerate(data):
    for j, text in enumerate([feat, plant, animal]):
        cell = table.cell(i + 1, j)
        cell.text = text
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.color.rgb = GRAY_LIGHT if j > 0 else WHITE
            p.font.bold = (j == 0)
            p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.RIGHT
        cell.fill.solid()
        cell.fill.fore_color.rgb = BG_CARD if i % 2 == 0 else BG_DARK

# Important note
note_card = add_shape_card(slide3, 0.5, 6.2, 12.333, 0.6)
add_textbox(slide3, 0.7, 6.25, 11.5, 0.5, 
            "نقطة مهمة: الخلية النباتية بتعمل تركيب ضوئي بسبب البلاستيدات الخضراء - شي الخلية الحيوانية ما بتعملو", 
            font_size=14, color=GRAY_LIGHT)

add_footer(slide3, 3)


# ===== SLIDE 4: Summary =====
slide4 = prs.slides.add_slide(blank_layout)
add_bg(slide4)

add_textbox(slide4, 0.5, 0.3, 6, 0.4, "مثال عملي - أدوات الذكاء الاصطناعي", 
            font_size=13, color=AMBER, bold=True)

add_textbox(slide4, 0.5, 0.9, 8, 0.6, "ملخص الدرس - النقاط الأساسية", 
            font_size=30, color=AMBER, bold=True)

add_textbox(slide4, 0.5, 1.5, 8, 0.4, "ملخص جاهز للتوزيع على الطلاب", 
            font_size=16, color=GRAY)

# 4 summary points
points = [
    ("1", "الخلية هي وحدة بناء الكائن الحي", "كل الكائنات الحية مكونة من خلايا", AMBER),
    ("2", "كل خلية فيها غشاء وسيتوبلازم ونواة", "هاد التركيب الأساسي المشترك", GREEN),
    ("3", "الخلية النباتية عندها مكونات إضافية", "جدار خلوي + بلاستيدات + فجوة كبيرة", PURPLE),
    ("4", "النواة هي مركز القيادة", "فيها DNA اللي بيحدد كل صفات الخلية", RED),
]

for i, (num, title, desc, color) in enumerate(points):
    y = 2.1 + i * 1.1
    add_shape_card(slide4, 0.5, y, 6, 0.9)
    
    # Number circle
    circle = slide4.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), Inches(y + 0.15), Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(0x1a, 0x1a, 0x1a)
    circle.line.color.rgb = color
    circle.line.width = Pt(2)
    tf = circle.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = color
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    add_textbox(slide4, 1.4, y + 0.05, 4.8, 0.4, title, font_size=15, color=WHITE, bold=True)
    add_textbox(slide4, 1.4, y + 0.45, 4.8, 0.4, desc, font_size=12, color=GRAY)

# Right side - Objectives card
add_shape_card(slide4, 7, 2.1, 5.8, 2.5)
add_textbox(slide4, 7.2, 2.2, 5.3, 0.4, "أهداف الدرس", font_size=18, color=AMBER, bold=True)
objectives = [
    "يتعرف الطالب على مكونات الخلية",
    "يقارن بين الخلية النباتية والحيوانية",
    "يفهم وظيفة كل مكون من مكونات الخلية",
    "يربط بين تركيب الخلية ووظيفتها",
]
for i, obj in enumerate(objectives):
    add_textbox(slide4, 7.4, 2.7 + i * 0.35, 5.1, 0.35, f"  {obj}", 
                font_size=13, color=GRAY_LIGHT)

# University & trainers section
add_shape_card(slide4, 7, 4.8, 5.8, 1.5)
add_textbox(slide4, 7.2, 4.9, 5.3, 0.4, "للدكاترة والمدربين", font_size=18, color=PURPLE, bold=True)
add_textbox(slide4, 7.4, 5.3, 5.1, 0.35, "  قابل للتطوير لمستوى جامعي", font_size=13, color=GRAY_LIGHT)
add_textbox(slide4, 7.4, 5.65, 5.1, 0.35, "  مناسب لورشات تدريب المعلمين", font_size=13, color=GRAY_LIGHT)

add_footer(slide4, 4)


# ===== SLIDE 5: Quiz Questions =====
slide5 = prs.slides.add_slide(blank_layout)
add_bg(slide5)

add_textbox(slide5, 0.5, 0.3, 6, 0.4, "مثال عملي - أدوات الذكاء الاصطناعي", 
            font_size=13, color=AMBER, bold=True)

add_textbox(slide5, 0.5, 0.9, 10, 0.6, "أسئلة تقييمية - اختبار سريع", 
            font_size=30, color=AMBER, bold=True)

add_textbox(slide5, 0.5, 1.5, 10, 0.4, "5 أسئلة متنوعة مع الإجابات النموذجية - جاهزة للطباعة", 
            font_size=16, color=GRAY)

questions = [
    ("1", "ما هي المكونات الأساسية المشتركة بين الخلية النباتية والحيوانية؟", "الغشاء الخلوي، السيتوبلازم، النواة (صفحة 42)"),
    ("2", "ما وظيفة النواة في الخلية؟", "تحتوي على الحمض النووي DNA وتتحكم بنشاطات الخلية وصفاتها (صفحة 44)"),
    ("3", "لماذا الخلية النباتية تعمل تركيب ضوئي والحيوانية لا؟", "لأنها تحتوي على البلاستيدات الخضراء التي تحتوي اليخضور (صفحة 49)"),
    ("4", "ما الفرق بين الفجوة العصارية في الخلية النباتية والحيوانية؟", "النباتية: كبيرة ومركزية واحدة. الحيوانية: صغيرة ومتعددة (صفحة 50)"),
    ("5", "قارن بين الجدار الخلوي والغشاء الخلوي من حيث الموقع والوظيفة؟", "الجدار: خارجي في النبات فقط، دعم وحماية. الغشاء: ينظم دخول وخروج المواد (صفحة 43)"),
]

for i, (num, q, a) in enumerate(questions):
    y = 2.0 + i * 1.05
    
    # Question box
    add_shape_card(slide5, 0.5, y, 12.333, 0.45, border_color=PURPLE)
    add_textbox(slide5, 0.7, y + 0.05, 0.4, 0.35, num, font_size=13, color=PURPLE, bold=True)
    add_textbox(slide5, 1.2, y + 0.05, 11.3, 0.35, q, font_size=13, color=WHITE, bold=True)
    
    # Answer
    add_textbox(slide5, 1.2, y + 0.48, 11.3, 0.4, a, font_size=12, color=GREEN)

# Note for teachers
note_card = add_shape_card(slide5, 0.5, 6.3, 12.333, 0.5, border_color=AMBER)
add_textbox(slide5, 0.7, 6.35, 11.5, 0.4, 
            "ملاحظة: هالأسئلة تم إعدادها من الكتاب المدرسي فقط - كل إجابة مكتوب جنبها رقم الصفحة. فيك تطلب نسخة تانية بأسئلة مختلفة بنقرة واحدة!", 
            font_size=12, color=GRAY_LIGHT)

add_footer(slide5, 5)


# ===== SLIDE 6: Lesson Plan =====
slide6 = prs.slides.add_slide(blank_layout)
add_bg(slide6)

add_textbox(slide6, 0.5, 0.3, 6, 0.4, "مثال عملي - أدوات الذكاء الاصطناعي", 
            font_size=13, color=AMBER, bold=True)

add_textbox(slide6, 0.5, 0.9, 8, 0.6, "خطة الدرس - تحضير كامل جاهز", 
            font_size=30, color=AMBER, bold=True)

add_textbox(slide6, 0.5, 1.5, 8, 0.4, "خطة درس منسقة ومكتملة - جاهزة للتقديم", 
            font_size=16, color=GRAY)

# Steps on the left
steps = [
    ("1", "التمهيد (5 دقائق)", 'سؤال تحفيزي: "شو أصغر شي حي بالعالم؟" - ربط بالتجربة اليومية'),
    ("2", "الشرح (15 دقيقة)", "عرض مكونات الخلية (غشاء-سيتوبلازم-نواة) مع الرسم التوضيحي"),
    ("3", "المقارنة (10 دقائق)", "جدول مقارنة بين الخلية النباتية والحيوانية - عمل الطلاب"),
    ("4", "التقييم (10 دقائق)", "5 أسئلة تقييمية متنوعة مع التصحيح الفوري"),
    ("5", "الواجب المنزلي", "رسم خلية نباتية مع تسمية المكونات + سؤال بحثي"),
]

for i, (num, title, desc) in enumerate(steps):
    y = 2.0 + i * 0.95
    
    circle = slide6.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), Inches(y), Inches(0.4), Inches(0.4))
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(0x1a, 0x1a, 0x1a)
    circle.line.color.rgb = AMBER
    circle.line.width = Pt(2)
    tf = circle.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = AMBER
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    add_textbox(slide6, 1.3, y - 0.05, 5.5, 0.35, title, font_size=15, color=WHITE, bold=True)
    add_textbox(slide6, 1.3, y + 0.3, 5.5, 0.35, desc, font_size=12, color=GRAY)

# Right side - Info cards
add_shape_card(slide6, 7, 2.0, 5.8, 1.5)
add_textbox(slide6, 7.2, 2.1, 5.3, 0.4, "بيانات الدرس", font_size=16, color=AMBER, bold=True)

info_items = [("المادة:", "علوم"), ("الصف:", "السابع"), ("المدة:", "45 دقيقة"), ("المرجع:", "ص 42-50")]
for i, (label, value) in enumerate(info_items):
    col = i % 2
    row = i // 2
    x = 7.4 + col * 2.7
    y = 2.6 + row * 0.4
    add_textbox(slide6, x, y, 1.2, 0.35, label, font_size=12, color=GRAY)
    add_textbox(slide6, x + 1.0, y, 1.5, 0.35, value, font_size=13, color=WHITE, bold=True)

add_shape_card(slide6, 7, 3.7, 5.8, 1.5)
add_textbox(slide6, 7.2, 3.8, 5.3, 0.4, "الأهداف", font_size=16, color=GREEN, bold=True)
obj_list = [
    "أن يتعرف الطالب على مكونات الخلية",
    "أن يقارن بين نوعي الخلايا",
    "أن يربط بين التركيب والوظيفة",
]
for i, obj in enumerate(obj_list):
    add_textbox(slide6, 7.4, 4.3 + i * 0.3, 5.1, 0.3, f"  {obj}", font_size=12, color=GRAY_LIGHT)

# Green highlight card
add_shape_card(slide6, 7, 5.4, 5.8, 1.2, border_color=GREEN)
add_textbox(slide6, 7.2, 5.5, 5.3, 0.4, "كل هذا تم إعداده بـ 5 دقائق!", 
            font_size=16, color=GREEN, bold=True)
add_textbox(slide6, 7.2, 5.9, 5.3, 0.6, 
            "بدل ساعتين تحضير، فيك تعمل خطة درس كاملة + أسئلة + ملخص + ملف PDF جاهز للطباعة. هاد اللي رح تتعلمه بالكورس.", 
            font_size=12, color=GRAY_LIGHT)

add_footer(slide6, 6)


# Save
prs.save(OUTPUT_PATH)
print(f"PPTX created: {OUTPUT_PATH}")
